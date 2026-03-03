"""Document Parser — handles any file type for RAG ingestion."""
import os
import io
import re
import logging
from typing import Any, Dict, List, Tuple

logger = logging.getLogger("mentora.parser")

# PDF
import fitz  # PyMuPDF

# OCR
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# DOCX
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# CSV / Excel
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# HTML / XML
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

# RTF
try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

# EPUB
try:
    import ebooklib
    from ebooklib import epub
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False


# ─── Text Utilities ──────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """Normalize whitespace and remove common noise characters."""
    text = re.sub(r"[\x00\x0c\r]", "", text)  # remove null/form-feed/CR
    text = re.sub(r"[ \t]+", " ", text)          # collapse horizontal whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)        # max two consecutive newlines
    return text.strip()


# Sentence boundary: period/!/? followed by whitespace then a capital, digit, or quote
_SENT_SPLIT = re.compile(r'(?<=[.!?])\s+(?=[A-Z\d"\'\'(])')


def chunk_text(text: str, chunk_size: int = 200, overlap: int = 40) -> List[str]:
    """
    Sentence-aware overlapping chunker.

    Splits at sentence boundaries first so embeddings capture complete thoughts.
    Target ≈ 200 words per chunk (≈ 260 tokens), well inside
    mxbai-embed-large's 512-token limit.  Overlap carries context across
    chunk boundaries so retrieval doesn't miss split answers.
    """
    text = _clean_text(text)
    sentences = [s.strip() for s in _SENT_SPLIT.split(text) if s.strip()]
    if not sentences:
        sentences = [text]

    chunks: List[str] = []
    buf: List[str] = []
    buf_words = 0

    for sent in sentences:
        words = sent.split()
        word_count = len(words)

        # Single sentence longer than chunk_size → sub-split at word boundary
        if word_count > chunk_size:
            if buf:
                chunks.append(" ".join(buf))
                buf = buf[-overlap:]
                buf_words = len(buf)
            # Sub-split the long sentence
            step = chunk_size - overlap
            for i in range(0, len(words), step):
                sub = words[i: i + chunk_size]
                if sub:
                    chunks.append(" ".join(sub))
            # Keep tail for overlap into next sentence
            buf = words[-overlap:] if len(words) >= overlap else words[:]
            buf_words = len(buf)
            continue

        if buf_words + word_count > chunk_size and buf:
            chunks.append(" ".join(buf))
            buf = buf[-overlap:]   # overlap carried forward
            buf_words = len(buf)

        buf.extend(words)
        buf_words += word_count

    if buf:
        chunks.append(" ".join(buf))

    return [c for c in chunks if c.strip()] or [text.strip()]


def _pdf_pages(file_path: str) -> List[Tuple[str, int]]:
    """
    Return a list of (clean_page_text, 1-based_page_number) for every
    non-empty page in a PDF.  Falls back to OCR on image-only pages.
    """
    doc = fitz.open(file_path)
    pages: List[Tuple[str, int]] = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if not text and OCR_AVAILABLE:
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img).strip()
        if text:
            pages.append((_clean_text(text), page_num))
    return pages


def parse_pdf(file_path: str) -> Tuple[str, int]:
    """Return (full_text, page_count). Used for summary; chunking uses _pdf_pages."""
    pages = _pdf_pages(file_path)
    return "\n\n".join(p[0] for p in pages), len(pages)


def parse_pptx(file_path: str) -> Tuple[str, int]:
    if not PPTX_AVAILABLE:
        return "", 0
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            slides_text.append("\n".join(slide_parts))
    return "\n\n".join(slides_text), len(slides_text)


def parse_docx(file_path: str) -> Tuple[str, int]:
    if not DOCX_AVAILABLE:
        return "", 0
    d = docx.Document(file_path)
    paragraphs = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), len(paragraphs)


def parse_txt(file_path: str) -> Tuple[str, int]:
    try:
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        text = raw.decode(enc, errors="ignore")
    except Exception:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            text = f.read()
    lines = [l for l in text.splitlines() if l.strip()]
    return text, len(lines)


def parse_image(file_path: str) -> Tuple[str, int]:
    if not OCR_AVAILABLE:
        return "", 0
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return text.strip(), 1


def parse_csv(file_path: str) -> Tuple[str, int]:
    if not PANDAS_AVAILABLE:
        return "", 0
    try:
        df = pd.read_csv(file_path, encoding_errors="ignore")
        text = df.to_string(index=False)
        return text, len(df)
    except Exception:
        return "", 0


def parse_excel(file_path: str) -> Tuple[str, int]:
    if not PANDAS_AVAILABLE:
        return "", 0
    try:
        sheets = pd.read_excel(file_path, sheet_name=None)
        parts = []
        total_rows = 0
        for name, df in sheets.items():
            parts.append(f"[Sheet: {name}]\n{df.to_string(index=False)}")
            total_rows += len(df)
        return "\n\n".join(parts), total_rows
    except Exception:
        return "", 0


def parse_html(file_path: str) -> Tuple[str, int]:
    if not BS4_AVAILABLE:
        with open(file_path, encoding="utf-8", errors="ignore") as f:
            return f.read(), 1
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    lines = [l for l in text.splitlines() if l.strip()]
    return "\n".join(lines), len(lines)


def parse_rtf(file_path: str) -> Tuple[str, int]:
    if not RTF_AVAILABLE:
        return "", 0
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        content = f.read()
    text = rtf_to_text(content)
    lines = [l for l in text.splitlines() if l.strip()]
    return text, len(lines)


def parse_epub(file_path: str) -> Tuple[str, int]:
    if not EPUB_AVAILABLE or not BS4_AVAILABLE:
        return "", 0
    book = epub.read_epub(file_path)
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if text.strip():
            parts.append(text)
    text = "\n\n".join(parts)
    return text, len(parts)


def parse_json(file_path: str) -> Tuple[str, int]:
    import json
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return text, 1


def parse_markdown(file_path: str) -> Tuple[str, int]:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        text = f.read()
    # Strip markdown syntax for cleaner embedding
    text = re.sub(r"#{1,6}\s*", "", text)              # headings
    text = re.sub(r"\*\*|__|\*|_|`{1,3}", "", text)   # bold/italic/code
    text = re.sub(r"!?\[.*?\]\(.*?\)", "", text)        # links/images
    text = re.sub(r"^[-*+>]\s+", "", text, flags=re.M) # lists/blockquotes
    lines = [l for l in text.splitlines() if l.strip()]
    return "\n".join(lines), len(lines)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def parse_file(file_path: str) -> Tuple[str, int, List[Dict[str, Any]]]:
    """
    Parse any file and return (full_text, page_count, chunks).

    Each chunk is a dict::

        {
            "text":         str,   # chunk content (≤ 200 words)
            "page":         int,   # 1-based page/slide number (0 = unknown)
            "chunk_index":  int,   # sequential index within this document
            "source_type":  str,   # file extension without leading dot
        }

    PDF and PPTX are chunked section-by-section to preserve page numbers.
    All other formats produce page=0 chunks.
    """
    ext = os.path.splitext(file_path)[-1].lower()
    src = ext.lstrip(".")
    idx = 0
    logger.info("📄 [PARSER] Parsing '%s'  type=%s", os.path.basename(file_path), ext or "unknown")

    # ── PDF: chunk per page so page numbers are preserved ──────────
    if ext == ".pdf":
        page_list = _pdf_pages(file_path)
        all_text_parts: List[str] = []
        chunks: List[Dict[str, Any]] = []
        for page_text, page_num in page_list:
            all_text_parts.append(page_text)
            for c in chunk_text(page_text):
                chunks.append({"text": c, "page": page_num,
                                "chunk_index": idx, "source_type": src})
                idx += 1
        logger.info("✅ [PARSER] PDF — %d pages → %d chunks", len(page_list), len(chunks))
        return "\n\n".join(all_text_parts), len(page_list), chunks

    # ── PPTX: chunk per slide ────────────────────────────────────────
    if ext in (".pptx", ".ppt") and PPTX_AVAILABLE:
        prs = Presentation(file_path)
        all_text_parts = []
        chunks = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts = [shape.text.strip() for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                slide_text = _clean_text("\n".join(parts))
                all_text_parts.append(slide_text)
                for c in chunk_text(slide_text):
                    chunks.append({"text": c, "page": slide_num,
                                   "chunk_index": idx, "source_type": src})
                    idx += 1
        logger.info("✅ [PARSER] PPTX — %d slides → %d chunks", len(prs.slides), len(chunks))
        return "\n\n".join(all_text_parts), len(prs.slides), chunks

    # ── All other formats: flat text then chunk ──────────────────────
    if ext in (".docx", ".doc"):
        text, pages = parse_docx(file_path)
    elif ext in (".xlsx", ".xls", ".xlsm"):
        text, pages = parse_excel(file_path)
    elif ext == ".csv":
        text, pages = parse_csv(file_path)
    elif ext in (".html", ".htm", ".xml"):
        text, pages = parse_html(file_path)
    elif ext == ".rtf":
        text, pages = parse_rtf(file_path)
    elif ext == ".epub":
        text, pages = parse_epub(file_path)
    elif ext == ".json":
        text, pages = parse_json(file_path)
    elif ext in (".md", ".markdown"):
        text, pages = parse_markdown(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp", ".gif"):
        text, pages = parse_image(file_path)
    else:
        text, pages = parse_txt(file_path)

    text = _clean_text(text)
    raw = chunk_text(text) if text else []
    chunks = [{"text": c, "page": 0, "chunk_index": i, "source_type": src}
              for i, c in enumerate(raw) if c.strip()]
    logger.info("✅ [PARSER] %s — %d pages → %d chunks", ext, pages, len(chunks))
    return text, pages, chunks
